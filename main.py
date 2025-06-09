import os
import json
import re
import fitz  # PyMuPDF for PDF text extraction
import docx  # python-docx for DOCX extraction
from pptx import Presentation
from pathlib import Path
from typing import List, Optional, Union

from fastapi import FastAPI, UploadFile, Form, File
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from dotenv import load_dotenv
import google.generativeai as genai

# Load .env file and API key
load_dotenv()
API_KEY = os.getenv("gemniKey")
if not API_KEY:
    raise ValueError("Missing GEMINI_API_KEY in .env")

genai.configure(api_key=API_KEY)
model = genai.GenerativeModel("gemini-1.5-flash")

UPLOAD_DIR = Path("Uploads")
UPLOAD_DIR.mkdir(exist_ok=True)

app = FastAPI()
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

class McqQuestion(BaseModel):
    questionNumber: int
    question: str
    options: List[str]
    correctAnswer: str
    explanation: str

def extract_text_from_file(file_path: str) -> str:
    ext = file_path.split('.')[-1].lower()
    try:
        if ext == "pdf":
            doc = fitz.open(file_path)
            return "\n".join(page.get_text("text") for page in doc)
        elif ext == "docx":
            doc = docx.Document(file_path)
            return "\n".join(para.text for para in doc.paragraphs)
        elif ext == "txt":
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        elif ext == "pptx":
            prs = Presentation(file_path)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs)
        else:
            return ""
    except Exception as e:
        print("Error reading file:", e)
        return ""

# ... (keep all your imports and setup unchanged)

@app.post("/generateQuestions")
async def generate_questions(
    sourceType: str = Form(...),
    textInput: Optional[str] = Form(None),
    numOfQuestions: int = Form(...),
    difficultyLevel: str = Form(...),
    typeOfQuestions: str = Form(...),
    language: str = Form("English"),
    fileInput: Optional[UploadFile] = File(None)
) -> Union[List[McqQuestion], dict]:

    if sourceType == "TEXT":
        text_source = textInput
    elif sourceType == "FILE" and fileInput:
        save_to = UPLOAD_DIR / fileInput.filename
        with open(save_to, 'wb') as f:
            f.write(await fileInput.read())
        text_source = extract_text_from_file(str(save_to))
        os.remove(save_to)
        if not text_source.strip():
            return {"error": "Could not extract text from the file. Ensure it's a readable PDF, DOCX, TXT or PPTX."}
    else:
        return {"error": "Invalid input. Provide either text or a valid file."}

    prompt = f"""
Generate exactly {numOfQuestions} {difficultyLevel} {typeOfQuestions} questions in {language} based on this text:

{text_source}

Each question must have these fields exactly:
- questionNumber: integer
- question: string
- options: list of strings
- correctAnswer: string
- explanation: string (a detailed explanation supporting the correct answer, cannot be empty)

Output ONLY a valid JSON array of objects with these fields, no markdown, no extra text.

If the language is Arabic, translate *all fields* including the explanation into Arabic fully and provide detailed explanations that are NOT empty.

Example format:

[
  {{
    "questionNumber": 1,
    "question": "ما هي عاصمة فرنسا؟",
    "options": ["باريس", "لندن", "برلين", "مدريد"],
    "correctAnswer": "باريس",
    "explanation": "باريس هي عاصمة فرنسا لأنها مقر الحكومة الفرنسية."
  }}
]
"""

    try:
        response = model.generate_content([prompt])
        response_text = response.text

        print("🧠 Raw Gemini Response:\n", response_text)

        cleaned_response = re.sub(r"^```json|```$", "", response_text.strip(), flags=re.MULTILINE).strip()
        print("🧠 Cleaned Response:\n", cleaned_response)

        parsed_data = json.loads(cleaned_response)

        # Fill fallback explanation if missing or empty
        questions = []
        for q in parsed_data:
            if not q.get("explanation"):
                q["explanation"] = "Explanation not provided by the AI."
            questions.append(McqQuestion(**q))

        return {"questionData": questions}

    except json.JSONDecodeError:
        return {"error": "The model did not return valid JSON. Try reducing input size or rephrasing text."}
    except Exception as e:
        print("Error:", e)
        return {"error": f"Gemini API failed: {str(e)}"}
